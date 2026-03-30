"""
lo.core.pptx_ops — Surgical shape operations on a Presentation.

All operations:
  - Are surgical: only touch the target shape
  - Use px units (converted to EMU internally, PX = 9525)
  - Save immediately after each operation
"""

from __future__ import annotations

import os
import re

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.util import Emu

from .shape_finder import find_shape_by_name

PX = 9525  # EMU per pixel at 96 DPI
HOST = "localhost"
PORT = 2002


def _auto_reload(pptx_path: str) -> None:
    """Attempt to hot-reload PPTX in running LibreOffice via UNO."""
    try:
        import uno
        from com.sun.star.beans import PropertyValue
    except ImportError:
        return  # python-uno not available, skip

    try:
        localContext = uno.getComponentContext()
        resolver = localContext.ServiceManager.createInstanceWithContext(
            "com.sun.star.bridge.UnoUrlResolver", localContext
        )
        ctx = resolver.resolve(
            f"uno:socket,host={HOST},port={PORT};urp;StarOffice.ComponentContext"
        )
        smgr = ctx.ServiceManager
        desktop = smgr.createInstanceWithContext("com.sun.star.frame.Desktop", ctx)

        abs_path = os.path.abspath(pptx_path)
        url = uno.systemPathToFileUrl(abs_path)

        # Close existing document if open
        for comp in desktop.Components:
            try:
                if hasattr(comp, "getURL") and comp.getURL() == url:
                    comp.close(True)
                    break
            except Exception:
                pass

        # Reopen
        props = []
        pv = PropertyValue()
        pv.Name = "MacroExecutionMode"
        pv.Value = 4
        props.append(pv)
        desktop.loadComponentFromURL(url, "_blank", 0, tuple(props))
        return  # success
    except Exception:
        pass  # LibreOffice not running with UNO socket, skip reload

    # Try launching LibreOffice with UNO socket (non-blocking)
    import subprocess
    cmd = [
        "soffice",
        f"--accept=socket,host={HOST},port={PORT};urp;StarOffice.ServiceManager",
        "--impress",
        os.path.abspath(pptx_path),
    ]
    subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


# ── Color ──────────────────────────────────────────────────────────────────────

_COLOR_HEX = re.compile(r"^#[0-9A-Fa-f]{6}$")


def _parse_color(color: str) -> int:
    """Parse "#RRGGBB" → int RGB."""
    if not _COLOR_HEX.match(color):
        raise ValueError(f"Invalid color format: {color!r}. Use #RRGGBB.")
    return int(color[1:], 16)


# ── Text ───────────────────────────────────────────────────────────────────────

def set_text(shape: BaseShape, text: str) -> None:
    """Replace all text in shape. Formatting must be re-applied after this call."""
    if not shape.has_text_frame:
        raise ValueError(f"Shape {shape.name!r} has no text frame")
    tf = shape.text_frame
    for para in list(tf.paragraphs)[1:]:
        para.clear()
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = text
        for run in list(first_para.runs)[1:]:
            run.text = ''
    else:
        first_para.text = text


def set_font_size(shape: BaseShape, size_pt: int) -> None:
    """Set font size (pt) for all runs in shape's text frame."""
    if not shape.has_text_frame:
        raise ValueError(f"Shape {shape.name!r} has no text frame")
    from pptx.util import Pt
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)


# ── Fill ──────────────────────────────────────────────────────────────────────

def set_fill(shape: BaseShape, color: str) -> None:
    """Set solid fill color. color is "#RRGGBB"."""
    rgb = _parse_color(color)
    shape.fill.solid()
    from pptx.dml.color import RGBColor
    shape.fill.fore_color.rgb = RGBColor(
        (rgb >> 16) & 0xFF,
        (rgb >> 8) & 0xFF,
        rgb & 0xFF,
    )


# ── Text Color ─────────────────────────────────────────────────────────────────

def set_text_color(shape: BaseShape, color: str) -> None:
    """Set the font color of the first text run. color is "#RRGGBB"."""
    if not shape.has_text_frame:
        raise ValueError(f"Shape {shape.name!r} has no text frame")
    rgb = _parse_color(color)
    from pptx.dml.color import RGBColor
    col = RGBColor((rgb >> 16) & 0xFF, (rgb >> 8) & 0xFF, rgb & 0xFF)
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.color.rgb = col


# ── Shape Info ─────────────────────────────────────────────────────────────────

def get_shape_info(shape: BaseShape) -> dict:
    """
    Return a dict of shape properties (no save).
    Includes: name, type, x, y, w, h, fill, text.
    """
    from pptx.dml.color import RGBColor
    info = {
        "name": shape.name,
        "type": shape.shape_type.name,
        "x": shape.left // PX,
        "y": shape.top // PX,
        "w": shape.width // PX,
        "h": shape.height // PX,
        "text": shape.text_frame.text if shape.has_text_frame else "",
    }
    try:
        fill = shape.fill
        if fill.type is not None and fill.type != 1:  # 1 = background/no fill
            info["fill"] = f"#{fill.fore_color.rgb}"
        else:
            info["fill"] = None
    except Exception:
        info["fill"] = None
    return info


# ── Position & Size ───────────────────────────────────────────────────────────

def move(shape: BaseShape, x_px: int, y_px: int) -> None:
    """Move shape to absolute position (px)."""
    shape.left = Emu(x_px * PX)
    shape.top = Emu(y_px * PX)


def resize(shape: BaseShape, w_px: int, h_px: int) -> None:
    """Resize shape to absolute dimensions (px)."""
    shape.width = Emu(w_px * PX)
    shape.height = Emu(h_px * PX)


# ── Apply + Save ───────────────────────────────────────────────────────────────

def apply_operation(
    prs: Presentation,
    pptx_path: str,
    name: str,
    op: str,
    value,
) -> None:
    """
    Find shape by name, apply one operation, then save.

    Operations:
      set-text   → value is str
      set-fill   → value is "#RRGGBB"
      move       → value is (x_px, y_px)
      resize     → value is (w_px, h_px)
    """
    shape = find_shape_by_name(prs, name)
    if shape is None:
        raise KeyError(f"Shape not found: {name!r}")

    if op == "set-text":
        set_text(shape, value)
    elif op == "set-fill":
        set_fill(shape, value)
    elif op == "move":
        x, y = value
        move(shape, x, y)
    elif op == "resize":
        w, h = value
        resize(shape, w, h)
    elif op == "set-text-color":
        set_text_color(shape, value)
    elif op == "set-font-size":
        set_font_size(shape, value)
    else:
        raise ValueError(f"Unknown operation: {op!r}")

    prs.save(pptx_path)

    # Auto-reload in LibreOffice if running
    _auto_reload(pptx_path)
