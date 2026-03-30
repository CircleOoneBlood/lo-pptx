"""
lo.core.shape_finder — Find a shape by name across all slides.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.shapes.base import BaseShape


def find_shape_by_name(prs: Presentation, name: str) -> BaseShape | None:
    """
    Search all slides for a shape with exact name match.
    Returns None if not found.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == name:
                return shape
    return None


def find_all_shapes_by_prefix(prs: Presentation, prefix: str) -> list[BaseShape]:
    """
    Find all shapes whose name starts with prefix.
    Useful for bulk operations like "s2_*" or "cover_*".
    """
    return [s for slide in prs.slides for s in slide.shapes if s.name.startswith(prefix)]
