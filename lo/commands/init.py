#!/usr/bin/env python3
"""
人类×AI 视觉协作工作流 — 小红书图文生成器
6 slides: 封面 + 5 张图文
画布: 1080×1350px (4:5 竖版)
"""

import json
import os
import shutil
import click

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── 画布 ──────────────────────────────────────────────────────────────────────
PX = 9525           # 1px in EMU at 96 DPI
W  = 1080 * PX
H  = 1350 * PX

# ── 色板 ──────────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0D, 0x0F, 0x1A)   # 深夜蓝
ACCENT      = RGBColor(0xFF, 0x3B, 0x5C)   # 小红书红
ACCENT2     = RGBColor(0x4E, 0xCD, 0xC4)   # 科技青
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0x88, 0x92, 0xA4)
TAG_BG      = RGBColor(0x1A, 0x20, 0x35)   # 标签底色
LINE        = RGBColor(0x2A, 0x31, 0x4A)   # 分割线

# ── 字体 ──────────────────────────────────────────────────────────────────────
FONT = "WenQuanYi Zen Hei"


# ── 基础绘图函数 ───────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, alpha=None, name=""):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                   Emu(x), Emu(y), Emu(w), Emu(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if name:
        shape.name = name
    return shape


def text(slide, txt, x, y, w, h, size, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False, name=""):
    txb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = txt
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = FONT
    if name:
        txb.name = name
    return txb


def multiline(slide, lines, x, y, w, h, size, bold=False,
              color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.3, name=""):
    """多行文本，lines 是 list of (text, color) 或 str"""
    txb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, col = item, color
        else:
            txt, col = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text           = txt
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = col
        run.font.name      = FONT
    if name:
        txb.name = name
    return txb


def bg(slide, sid):
    """全背景"""
    rect(slide, 0, 0, W, H, BG, name=f"{sid}_bg")


def top_accent_bar(slide, sid, h_px=8):
    rect(slide, 0, 0, W, h_px * PX, ACCENT, name=f"{sid}_accent_bar")


def bottom_bar(slide, sid, label="", h_px=80):
    by = H - h_px * PX
    rect(slide, 0, by, W, h_px * PX, TAG_BG, name=f"{sid}_bottom_bar")
    if label:
        text(slide, label, 60*PX, by + 18*PX, W - 120*PX, 50*PX,
             size=18, color=MUTED, align=PP_ALIGN.CENTER, name=f"{sid}_bottom_label")


def slide_number(slide, n, total, sid):
    label = f"{n:02d} / {total:02d}"
    text(slide, label, W - 180*PX, H - 70*PX, 140*PX, 40*PX,
         size=18, color=MUTED, align=PP_ALIGN.RIGHT, name=f"{sid}_slide_num")


def section_tag(slide, label, sid, x_px=60, y_px=60):
    """小标签块"""
    tw = (len(label) * 18 + 40) * PX
    rect(slide, x_px*PX, y_px*PX, tw, 44*PX, TAG_BG, name=f"{sid}_tag_bg")
    text(slide, label, (x_px+14)*PX, (y_px+8)*PX, tw - 20*PX, 30*PX,
         size=16, color=ACCENT, bold=True, name=f"{sid}_tag_text")


def divider(slide, y_px, sid, suffix="1", w_px=960):
    rect(slide, 60*PX, y_px*PX, w_px*PX, 2*PX, LINE, name=f"{sid}_divider_{suffix}")


def accent_dot(slide, x_px, y_px, sid, suffix="", r_px=10, color=ACCENT):
    rect(slide, x_px*PX, y_px*PX, r_px*2*PX, r_px*2*PX, color,
         name=f"{sid}_dot_{suffix}" if suffix else f"{sid}_dot")


# ── Slide 1: 封面 ──────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sid = "cover"
    bg(slide, sid)
    top_accent_bar(slide, sid, 6)

    # 顶部标签
    section_tag(slide, "# AI × 视觉设计", sid, 60, 50)

    # 大标题
    text(slide,
         "别跟 AI 描述了",
         60*PX, 180*PX, 960*PX, 160*PX,
         size=68, bold=True, color=WHITE, name="cover_title_main")
    text(slide,
         "直接拖给它看",
         60*PX, 330*PX, 960*PX, 160*PX,
         size=68, bold=True, color=ACCENT, name="cover_title_accent")

    # 分隔线
    rect(slide, 60*PX, 510*PX, 80*PX, 6*PX, ACCENT, name="cover_divider")

    # 副标题
    text(slide,
         "人类×AI 视觉内容协作工作流",
         60*PX, 540*PX, 960*PX, 80*PX,
         size=28, color=MUTED, name="cover_subtitle")

    # 三个关键词卡片
    keywords = ["拖拽编辑", "自动 Diff", "意图理解"]
    colors   = [ACCENT, ACCENT2, WHITE]
    for i, (kw, col) in enumerate(zip(keywords, colors)):
        kx = (60 + i * 330) * PX
        ky = 900 * PX
        rect(slide, kx, ky, 300*PX, 80*PX, TAG_BG, name=f"cover_kw_bg_{i}")
        text(slide, kw, kx + 20*PX, ky + 18*PX, 260*PX, 50*PX,
             size=20, bold=True, color=col, align=PP_ALIGN.CENTER,
             name=f"cover_kw_text_{i}")

    # 向下滚动提示（填补卡片和底部之间）
    divider(slide, 1010, sid, "scroll", w_px=840)
    text(slide, "↓  打开 PPT 直接动手改", 60*PX, 1030*PX, 960*PX, 50*PX,
         size=18, color=MUTED, align=PP_ALIGN.CENTER, name="cover_scroll_hint")
    rect(slide, 420*PX, 1100*PX, 120*PX, 4*PX, ACCENT, name="cover_scroll_accent")

    # 底部说明
    text(slide,
         "视觉细节用操作表达，不用语言描述",
         60*PX, 1160*PX, 960*PX, 60*PX,
         size=22, color=MUTED, name="cover_caption")

    # 装饰：右下角大数字
    text(slide, "01", 820*PX, 1160*PX, 200*PX, 120*PX,
         size=120, bold=True, color=LINE, align=PP_ALIGN.RIGHT, name="cover_deco_num")

    bottom_bar(slide, sid, "#ClaudeCode  #AI设计  #人机协作  #效率工具")


# ── Slide 2: 痛点 ──────────────────────────────────────────────────────────────

def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sid = "s2"
    bg(slide, sid)
    top_accent_bar(slide, sid)
    section_tag(slide, "痛点", sid, 60, 50)

    text(slide, "跟 AI 描述视觉细节",
         60*PX, 140*PX, 960*PX, 80*PX,
         size=40, bold=True, name="s2_headline")

    divider(slide, 240, sid, "1")

    # 场景描述
    problems = [
        "「把标题往左挪 20px」",
        "「颜色再深一点，不是这个深」",
        "说了十轮还没对齐。",
    ]
    y = 270
    for i, line in enumerate(problems):
        col = ACCENT if i == 2 else WHITE
        sz  = 38 if i == 2 else 30
        text(slide, line, 60*PX, y*PX, 960*PX, 70*PX,
             size=sz, bold=(i == 2), color=col, name=f"s2_problem_{i}")
        y += 80

    # 引号装饰
    text(slide, "❝", 850*PX, 260*PX, 180*PX, 120*PX,
         size=100, color=TAG_BG, align=PP_ALIGN.RIGHT, name="s2_quote_deco")

    # 痛点清单
    divider(slide, 580, sid, "2")
    text(slide, "用自然语言描述视觉意图的代价：",
         60*PX, 600*PX, 960*PX, 50*PX,
         size=22, color=MUTED, name="s2_list_intro")

    items = ["位置、大小无法精确表达", "颜色沟通来回试错", "多个元素的相对关系说不清", "一张图改十轮还不满意"]
    for i, item in enumerate(items):
        iy = 660 + i * 70
        rect(slide, 60*PX, iy*PX, 8*PX, 44*PX, ACCENT, name=f"s2_item_bar_{i}")
        text(slide, item, 90*PX, iy*PX, 900*PX, 50*PX,
             size=26, color=WHITE, name=f"s2_item_text_{i}")

    text(slide, "换个思路：让人类直接动手改 →",
         60*PX, 1180*PX, 960*PX, 50*PX,
         size=20, color=ACCENT2, name="s2_cta")

    slide_number(slide, 2, 6, sid)
    bottom_bar(slide, sid)


# ── Slide 3: 解法全貌 ──────────────────────────────────────────────────────────

def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sid = "s3"
    bg(slide, sid)
    top_accent_bar(slide, sid)
    section_tag(slide, "解法", sid, 60, 50)

    text(slide, "PPTX 当画布，人机各司其职",
         60*PX, 140*PX, 960*PX, 80*PX,
         size=38, bold=True, name="s3_headline")

    divider(slide, 245, sid, "1")

    # 流程图：三个步骤
    steps = [
        ("AI 生成初稿",   "lo init",     "一键生成完整 PPTX 模板", ACCENT),
        ("人类直接编辑",   "拖拽 / 改色",  "在 PPT 里拖到想要的位置", ACCENT2),
        ("AI 读懂改动",   "lo diff",     "自动对比变化，理解意图",  WHITE),
    ]

    for i, (hook, when, what, col) in enumerate(steps):
        base_y = 290 + i * 290

        rect(slide, 60*PX, base_y*PX, 6*PX, 240*PX, col,
             name=f"s3_step_bar_{i}")
        text(slide, hook,
             90*PX, base_y*PX, 500*PX, 70*PX,
             size=32, bold=True, color=col, name=f"s3_step_hook_{i}")
        text(slide, f"操作方式：{when}",
             90*PX, (base_y+70)*PX, 900*PX, 50*PX,
             size=20, color=MUTED, name=f"s3_step_when_{i}")
        text(slide, what,
             90*PX, (base_y+120)*PX, 900*PX, 60*PX,
             size=24, color=WHITE, name=f"s3_step_what_{i}")

        if i < 2:
            text(slide, "↓",
                 500*PX, (base_y+240)*PX, 80*PX, 50*PX,
                 size=28, color=LINE, align=PP_ALIGN.CENTER,
                 name=f"s3_arrow_{i}")

    slide_number(slide, 3, 6, sid)
    bottom_bar(slide, sid)


# ── Slide 4: lo diff 演示 ─────────────────────────────────────────────────────

def slide_stop_hook(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sid = "s4"
    bg(slide, sid)
    top_accent_bar(slide, sid)
    section_tag(slide, "核心能力", sid, 60, 50)

    text(slide, "lo diff", 60*PX, 130*PX, 700*PX, 80*PX,
         size=52, bold=True, color=ACCENT, name="s4_title")
    text(slide, "AI 自动读懂你改了什么", 60*PX, 210*PX, 700*PX, 60*PX,
         size=30, color=MUTED, name="s4_subtitle")

    divider(slide, 290, sid, "1")

    text(slide, "你在 PPT 里改完，AI 自动看到：",
         60*PX, 310*PX, 960*PX, 50*PX,
         size=22, color=MUTED, name="s4_intro")

    flow = [
        ("位置", "标题从 x=60 移到了 x=120"),
        ("颜色", "#FF3B5C → #4ECDC4"),
        ("文字", "「初稿」改成了「终稿」"),
        ("大小", "字号从 28pt 调到 36pt"),
    ]

    for i, (action, detail) in enumerate(flow):
        fy = 370 + i * 155
        rect(slide, 60*PX, fy*PX, 60*PX, 60*PX, ACCENT,
             name=f"s4_flow_num_bg_{i}")
        text(slide, str(i+1), 60*PX, (fy+8)*PX, 60*PX, 44*PX,
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             name=f"s4_flow_num_text_{i}")
        text(slide, action, 140*PX, fy*PX, 160*PX, 60*PX,
             size=26, bold=True, color=ACCENT, name=f"s4_flow_action_{i}")
        text(slide, detail, 310*PX, (fy+4)*PX, 710*PX, 52*PX,
             size=22, color=WHITE, name=f"s4_flow_detail_{i}")

        if i < len(flow) - 1:
            text(slide, "↓", 80*PX, (fy+65)*PX, 30*PX, 40*PX,
                 size=22, color=LINE, align=PP_ALIGN.CENTER,
                 name=f"s4_arrow_{i}")

    rect(slide, 60*PX, 1140*PX, 960*PX, 80*PX, TAG_BG, name="s4_result_bg")
    text(slide, "零描述成本，改动即沟通",
         80*PX, 1158*PX, 900*PX, 50*PX,
         size=20, color=ACCENT2, name="s4_result_text")

    slide_number(slide, 4, 6, sid)
    bottom_bar(slide, sid)


# ── Slide 5: 完整工作流 ───────────────────────────────────────────────────────

def slide_session_start(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sid = "s5"
    bg(slide, sid)
    top_accent_bar(slide, sid)
    section_tag(slide, "完整流程", sid, 60, 50)

    text(slide, "一次完整的协作循环", 60*PX, 130*PX, 960*PX, 80*PX,
         size=40, bold=True, color=ACCENT2, name="s5_title")

    divider(slide, 240, sid, "1")

    steps = [
        ("lo init",   "AI 生成 6 页小红书初稿"),
        ("打开 PPT",   "在 LibreOffice / WPS 里编辑"),
        ("发消息",     "告诉 AI「我改完了」"),
        ("lo diff",   "AI 对比快照，报告所有变化"),
        ("确认",       "满意 → lo snapshot 更新基准"),
        ("继续迭代",   "AI 根据你的意图跟进调整"),
    ]

    for i, (step, desc) in enumerate(steps):
        ky = 270 + i * 145
        col = ACCENT if i % 2 == 0 else ACCENT2
        rect(slide, 60*PX, ky*PX, 8*PX, 56*PX, col,
             name=f"s5_knows_bar_{i}")
        text(slide, step, 90*PX, (ky+4)*PX, 240*PX, 50*PX,
             size=22, color=col, bold=True, name=f"s5_knows_check_{i}")
        text(slide, desc, 340*PX, (ky+4)*PX, 640*PX, 50*PX,
             size=22, color=WHITE, name=f"s5_knows_text_{i}")

        if i < len(steps) - 1:
            text(slide, "↓", 80*PX, (ky+68)*PX, 30*PX, 40*PX,
                 size=18, color=LINE, align=PP_ALIGN.CENTER,
                 name=f"s5_arrow_{i}")

    slide_number(slide, 5, 6, sid)
    bottom_bar(slide, sid)


# ── Slide 6: 核心理念 + CTA ──────────────────────────────────────────────────

def slide_compact_sessions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sid = "s6"
    bg(slide, sid)
    top_accent_bar(slide, sid)
    section_tag(slide, "为什么这样做", sid, 60, 50)

    text(slide, "视觉细节用操作表达",
         60*PX, 140*PX, 960*PX, 70*PX,
         size=42, bold=True, color=WHITE, name="s6_pc_title")
    text(slide, "不用语言描述",
         60*PX, 220*PX, 960*PX, 70*PX,
         size=42, bold=True, color=ACCENT, name="s6_pc_title2")

    divider(slide, 310, sid, "1")

    text(slide,
         "直接拖到想要的位置，比跟 AI 描述快一个数量级。",
         60*PX, 330*PX, 960*PX, 60*PX,
         size=22, color=MUTED, name="s6_pc_intro")

    # 对比
    comparisons = [
        ("传统方式", "「把标题字号改大一点，颜色换成\n偏青色，往右挪一些」× 10 轮", ACCENT),
        ("协作工作流", "直接改 → AI 看到 → 一轮搞定", ACCENT2),
    ]

    for i, (label, desc, col) in enumerate(comparisons):
        cy = 420 + i * 220
        rect(slide, 60*PX, cy*PX, 960*PX, 190*PX, TAG_BG,
             name=f"s6_cmp_bg_{i}")
        rect(slide, 60*PX, cy*PX, 8*PX, 190*PX, col,
             name=f"s6_cmp_bar_{i}")
        text(slide, label, 90*PX, (cy+15)*PX, 380*PX, 44*PX,
             size=22, bold=True, color=col, name=f"s6_cmp_label_{i}")
        text(slide, desc, 90*PX, (cy+70)*PX, 900*PX, 100*PX,
             size=20, color=WHITE, name=f"s6_cmp_desc_{i}")

    # CTA
    rect(slide, 60*PX, 900*PX, 960*PX, 100*PX, ACCENT, name="s6_cta_bg")
    text(slide, "PPTX 是画布，不是目的",
         80*PX, 910*PX, 920*PX, 40*PX,
         size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         name="s6_cta_title")
    text(slide, "最终产物是图片 — 小红书配图、海报、社交媒体素材",
         80*PX, 955*PX, 920*PX, 40*PX,
         size=18, color=WHITE, align=PP_ALIGN.CENTER,
         name="s6_cta_desc")

    # 三个命令
    cmds = ["lo init", "lo diff", "lo snapshot"]
    for i, cmd in enumerate(cmds):
        cx = (60 + i * 330) * PX
        cy = 1060 * PX
        rect(slide, cx, cy, 300*PX, 70*PX, TAG_BG, name=f"s6_cmd_bg_{i}")
        text(slide, cmd, cx + 20*PX, cy + 16*PX, 260*PX, 40*PX,
             size=20, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER,
             name=f"s6_cmd_text_{i}")

    slide_number(slide, 6, 6, sid)
    bottom_bar(slide, sid, "#ClaudeCode  #AI设计  #人机协作  #效率工具  #小红书")


# ── 主程序 ─────────────────────────────────────────────────────────────────────

def _default_output() -> str:
    """Read output filename from pptx-workflow.json, fallback to 'output.pptx'."""
    wf = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "pptx-workflow.json")
    if os.path.exists(wf):
        with open(wf) as f:
            d = json.load(f)
        return d.get("pptx", "output.pptx")
    return "output.pptx"



@click.command()
@click.option("--output", "-o", "output_path", default=None,
              help="Output PPTX path (default: from pptx-workflow.json)")
def init(output_path: str | None):
    """Generate a new PPTX template for the current topic."""
    if output_path is None:
        output_path = _default_output()

    base = os.path.splitext(output_path)[0]
    baseline_path = f"{base}.baseline.pptx"

    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    slide_cover(prs)
    slide_problem(prs)
    slide_overview(prs)
    slide_stop_hook(prs)
    slide_session_start(prs)
    slide_compact_sessions(prs)

    prs.save(output_path)
    click.echo(f"✓ {output_path}  ({len(prs.slides)} slides)")

    if not os.path.exists(baseline_path):
        shutil.copy2(output_path, baseline_path)
        click.echo(f"✓ baseline: {baseline_path}")
    else:
        click.echo("  baseline exists, skipping.")


if __name__ == "__main__":
    main()
