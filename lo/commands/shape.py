"""
lo.commands.shape — shape subcommand group.
"""

from __future__ import annotations

import click

from ..core.config import load_config
from ..core.pptx_ops import apply_operation


@click.group(name="shape")
def shape_group():
    """Shape manipulation commands."""
    pass


@shape_group.command(name="set-text")
@click.option("--name", required=True, help="Shape name (globally unique)")
@click.option("--text", required=True, help="New text content")
@click.option("--file", "pptx_path", default=None, help="PPTX file path")
def set_text(name: str, text: str, pptx_path: str | None):
    """Set the text content of a shape."""
    cfg = load_config(pptx_path)
    from pptx import Presentation
    prs = Presentation(cfg.pptx_path)
    try:
        apply_operation(prs, cfg.pptx_path, name, "set-text", text)
        click.echo(f"✓ {name}: text set to {text!r}")
    except KeyError as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)
    except Exception as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)


@shape_group.command(name="set-fill")
@click.option("--name", required=True, help="Shape name")
@click.option("--color", required=True, help="Fill color as #RRGGBB")
@click.option("--file", "pptx_path", default=None, help="PPTX file path")
def set_fill(name: str, color: str, pptx_path: str | None):
    """Set the solid fill color of a shape."""
    cfg = load_config(pptx_path)
    from pptx import Presentation
    prs = Presentation(cfg.pptx_path)
    try:
        apply_operation(prs, cfg.pptx_path, name, "set-fill", color)
        click.echo(f"✓ {name}: fill set to {color}")
    except KeyError as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)
    except Exception as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)


@shape_group.command(name="move")
@click.option("--name", required=True, help="Shape name")
@click.option("--x", required=True, type=int, help="X position in px")
@click.option("--y", required=True, type=int, help="Y position in px")
@click.option("--file", "pptx_path", default=None, help="PPTX file path")
def move(name: str, x: int, y: int, pptx_path: str | None):
    """Move a shape to absolute position (px)."""
    cfg = load_config(pptx_path)
    from pptx import Presentation
    prs = Presentation(cfg.pptx_path)
    try:
        apply_operation(prs, cfg.pptx_path, name, "move", (x, y))
        click.echo(f"✓ {name}: moved to ({x}, {y})")
    except KeyError as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)
    except Exception as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)


@shape_group.command(name="resize")
@click.option("--name", required=True, help="Shape name")
@click.option("--w", required=True, type=int, help="Width in px")
@click.option("--h", required=True, type=int, help="Height in px")
@click.option("--file", "pptx_path", default=None, help="PPTX file path")
def resize(name: str, w: int, h: int, pptx_path: str | None):
    """Resize a shape to absolute dimensions (px)."""
    cfg = load_config(pptx_path)
    from pptx import Presentation
    prs = Presentation(cfg.pptx_path)
    try:
        apply_operation(prs, cfg.pptx_path, name, "resize", (w, h))
        click.echo(f"✓ {name}: resized to {w}x{h}")
    except KeyError as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)
    except Exception as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)


@shape_group.command(name="set-text-color")
@click.option("--name", required=True, help="Shape name")
@click.option("--color", required=True, help="Font color as #RRGGBB")
@click.option("--file", "pptx_path", default=None, help="PPTX file path")
def set_text_color(name: str, color: str, pptx_path: str | None):
    """Set the font color of text in a shape."""
    cfg = load_config(pptx_path)
    from pptx import Presentation
    prs = Presentation(cfg.pptx_path)
    try:
        apply_operation(prs, cfg.pptx_path, name, "set-text-color", color)
        click.echo(f"✓ {name}: text color set to {color}")
    except KeyError as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)
    except Exception as e:
        click.echo(f"✗ {e}", err=True)
        raise SystemExit(1)


@shape_group.command(name="get-info")
@click.option("--name", required=True, help="Shape name")
@click.option("--file", "pptx_path", default=None, help="PPTX file path")
def get_info(name: str, pptx_path: str | None):
    """Print shape properties (name, type, x, y, w, h, fill, text). No save."""
    cfg = load_config(pptx_path)
    from pptx import Presentation
    prs = Presentation(cfg.pptx_path)
    from ..core.shape_finder import find_shape_by_name
    from ..core.pptx_ops import get_shape_info
    shape = find_shape_by_name(prs, name)
    if shape is None:
        click.echo(f"✗ Shape not found: {name!r}", err=True)
        raise SystemExit(1)
    info = get_shape_info(shape)
    for k, v in info.items():
        click.echo(f"  {k}: {v}")
